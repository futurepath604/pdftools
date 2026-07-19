import os
import sys
import json
import shutil
from typing import List
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles

# পাইথন পাথ এনভায়রনমেন্ট সেটআপ
current_dir = os.path.dirname(os.path.abspath(__file__))
parent_dir = os.path.dirname(current_dir)
if current_dir not in sys.path:
    sys.path.append(current_dir)
if parent_dir not in sys.path:
    sys.path.append(parent_dir)

# --- SAFE IMPORTS BLOCK ---
# কোনো মডিউলে ইন্টারনাল এরর থাকলেও যেন আপনার পুরো সার্ভার ডাউন বা ক্র্যাশ না হয়
try: from app.tools.compressor import compress_pdf_file
except Exception: compress_pdf_file = None

try: from app.tools.pdf_to_image import pdf_to_images
except Exception: pdf_to_images = None

try: from app.tools.image_to_pdf import images_to_pdf
except Exception: images_to_pdf = None

try: from app.tools.modify import modify_pdf_pages
except Exception: modify_pdf_pages = None

try: from app.tools.security import lock_pdf_file, unlock_pdf_file
except Exception: lock_pdf_file = unlock_pdf_file = None

try: from app.tools.ocr_engine import process_pdf_ocr
except Exception: process_pdf_ocr = None

try: from app.tools.rearrange_backend import rearrange_pdf_pages
except Exception: rearrange_pdf_pages = None


# --- DYNAMIC BACKUP FOR PDF MERGER ---
def merge_pdf_files(input_paths: list, output_path: str):
    try:
        try:
            from pypdf import PdfMerger
            merger = PdfMerger()
            for path in input_paths: merger.append(path)
            merger.write(output_path)
            merger.close()
        except ImportError:
            import pypdf
            if hasattr(pypdf, 'PdfFileMerger'):
                merger = pypdf.PdfFileMerger()
                for path in input_paths: merger.append(path)
                merger.write(output_path)
                merger.close()
            else:
                from PyPDF2 import PdfMerger
                merger = PdfMerger()
                for path in input_paths: merger.append(path)
                merger.write(output_path)
                merger.close()
    except Exception as e:
        raise Exception(f"Merge Library Error: {str(e)}")


app = FastAPI(title="Secure PDF Tools Ultimate API")

# Mount Static Files for Frontend UI
app.mount("/static", StaticFiles(directory="app/static"), name="static")

# --- HTML UI ROUTES ---
@app.get("/")
async def read_index(): return FileResponse("app/static/index.html")

@app.get("/compress")
async def read_compress(): return FileResponse("app/static/compress.html")

@app.get("/merge")
async def read_merge(): return FileResponse("app/static/merge.html")

@app.get("/pdf-to-image")
async def read_pdf_to_image(): return FileResponse("app/static/pdf-to-image.html")

@app.get("/image-to-pdf")
async def read_image_to_pdf(): return FileResponse("app/static/image-to-pdf.html")

@app.get("/split")
async def read_split(): return FileResponse("app/static/split.html")

@app.get("/rotate")
async def read_rotate(): return FileResponse("app/static/rotate.html")

@app.get("/delete-pages")
async def read_delete_pages(): return FileResponse("app/static/delete-pages.html")

@app.get("/lock")
async def read_lock(): return FileResponse("app/static/lock.html")

@app.get("/unlock")
async def read_unlock(): return FileResponse("app/static/unlock.html")

# আপনার স্ট্রাকচারে থাকা নতুন টুলের রাউট
@app.get("/ocr")
async def read_ocr(): return FileResponse("app/static/ocr.html")

@app.get("/rearrange")
async def read_rearrange(): return FileResponse("app/static/rearrange.html")

@app.get("/office")
async def read_office(): return FileResponse("app/static/office.html")


# --- CORE API ENDPOINTS ---

@app.post("/api/compress")
async def compress_pdf(file: UploadFile = File(...)):
    if not compress_pdf_file: raise HTTPException(status_code=501, detail="Compress module not loaded")
    input_path = f"temp_{file.filename}"
    output_path = f"compressed_{file.filename}"
    with open(input_path, "wb") as buffer: shutil.copyfileobj(file.file, buffer)
    try:
        compress_pdf_file(input_path, output_path)
        return FileResponse(output_path, media_type="application/pdf", filename=output_path)
    except Exception as e: raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(input_path): os.remove(input_path)

@app.post("/api/merge")
async def merge_pdfs(files: List[UploadFile] = File(...)):
    input_paths = []
    output_path = "merged_document.pdf"
    for file in files:
        path = f"temp_{file.filename}"
        with open(path, "wb") as buffer: shutil.copyfileobj(file.file, buffer)
        input_paths.append(path)
    try:
        merge_pdf_files(input_paths, output_path)
        return FileResponse(output_path, media_type="application/pdf", filename=output_path)
    except Exception as e: raise HTTPException(status_code=500, detail=str(e))
    finally:
        for path in input_paths:
            if os.path.exists(path): os.remove(path)

@app.post("/api/pdf-to-image")
async def pdf_to_img(file: UploadFile = File(...)):
    if not pdf_to_images: raise HTTPException(status_code=501, detail="Converter module not loaded")
    input_path = f"temp_{file.filename}"
    output_zip = f"images_{file.filename}.zip"
    with open(input_path, "wb") as buffer: shutil.copyfileobj(file.file, buffer)
    try:
        pdf_to_images(input_path, output_zip)
        return FileResponse(output_zip, media_type="application/zip", filename=output_zip)
    except Exception as e: raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(input_path): os.remove(input_path)

@app.post("/api/image-to-pdf")
async def img_to_pdf(files: List[UploadFile] = File(...)):
    if not images_to_pdf: raise HTTPException(status_code=501, detail="Converter module not loaded")
    input_paths = []
    output_path = "images_converted.pdf"
    for file in files:
        path = f"temp_{file.filename}"
        with open(path, "wb") as buffer: shutil.copyfileobj(file.file, buffer)
        input_paths.append(path)
    try:
        images_to_pdf(input_paths, output_path)
        return FileResponse(output_path, media_type="application/pdf", filename=output_path)
    except Exception as e: raise HTTPException(status_code=500, detail=str(e))
    finally:
        for path in input_paths:
            if os.path.exists(path): os.remove(path)

@app.post("/api/modify-pdf")
async def modify_pdf(file: UploadFile = File(...), mode: str = Form(...), params: str = Form(...)):
    if not modify_pdf_pages: raise HTTPException(status_code=501, detail="Modify module not loaded")
    input_path = f"temp_{file.filename}"
    output_path = f"modified_{file.filename}"
    with open(input_path, "wb") as buffer: shutil.copyfileobj(file.file, buffer)
    try:
        param_dict = json.loads(params)
        modify_pdf_pages(input_path, output_path, mode, param_dict)
        return FileResponse(output_path, media_type="application/pdf", filename=output_path)
    except Exception as e: raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(input_path): os.remove(input_path)

@app.post("/api/security-pdf")
async def security_pdf(file: UploadFile = File(...), mode: str = Form(...), password: str = Form(...)):
    if not lock_pdf_file: raise HTTPException(status_code=501, detail="Security module not loaded")
    input_path = f"temp_{file.filename}"
    output_path = f"secured_{file.filename}"
    with open(input_path, "wb") as buffer: shutil.copyfileobj(file.file, buffer)
    try:
        if mode == "lock": lock_pdf_file(input_path, output_path, password)
        elif mode == "unlock": unlock_pdf_file(input_path, output_path, password)
        else: raise HTTPException(status_code=400, detail="Invalid security mode")
        return FileResponse(output_path, media_type="application/pdf", filename=output_path)
    except Exception as e: raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(input_path): os.remove(input_path)


# --- NEWLY ADDED ADVANCED CONVERSION ENDPOINTS (Word, Excel, PPT, OCR, Rearrange) ---

@app.post("/api/ocr")
async def ocr_pdf(file: UploadFile = File(...)):
    if not process_pdf_ocr: raise HTTPException(status_code=501, detail="OCR Engine not loaded")
    input_path = f"temp_{file.filename}"
    output_path = f"ocr_{file.filename}"
    with open(input_path, "wb") as buffer: shutil.copyfileobj(file.file, buffer)
    try:
        process_pdf_ocr(input_path, output_path)
        return FileResponse(output_path, media_type="application/pdf", filename=output_path)
    except Exception as e: raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(input_path): os.remove(input_path)

@app.post("/api/rearrange")
async def rearrange_pdf(file: UploadFile = File(...), page_order: str = Form(...)):
    if not rearrange_pdf_pages: raise HTTPException(status_code=501, detail="Rearrange module not loaded")
    input_path = f"temp_{file.filename}"
    output_path = f"rearranged_{file.filename}"
    with open(input_path, "wb") as buffer: shutil.copyfileobj(file.file, buffer)
    try:
        order_list = json.loads(page_order) # Expecting array of page numbers e.g. [1, 3, 2]
        rearrange_pdf_pages(input_path, output_path, order_list)
        return FileResponse(output_path, media_type="application/pdf", filename=output_path)
    except Exception as e: raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(input_path): os.remove(input_path)

@app.post("/api/pdf-to-office")
async def pdf_to_office(file: UploadFile = File(...), target_format: str = Form(...)):
    """
    Handles PDF to Word (.docx), PDF to Excel (.xlsx), and PDF to PPT (.pptx)
    target_format should be one of: 'word', 'excel', 'ppt'
    """
    input_path = f"temp_{file.filename}"
    base_name = os.path.splitext(file.filename)[0]
    
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        if target_format.lower() == "word":
            from pdf2docx import Converter
            output_path = f"{base_name}.docx"
            cv = Converter(input_path)
            cv.convert(output_path)
            cv.close()
            return FileResponse(output_path, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename=output_path)
            
        elif target_format.lower() == "excel":
            import pdfplumber
            import pandas as pd
            output_path = f"{base_name}.xlsx"
            all_tables = []
            with pdfplumber.open(input_path) as pdf:
                for page in pdf.pages:
                    tables = page.extract_tables()
                    for table in tables:
                        all_tables.append(pd.DataFrame(table))
            if all_tables:
                with pd.ExcelWriter(output_path) as writer:
                    for idx, df in enumerate(all_tables):
                        df.to_excel(writer, sheet_name=f"Table_{idx+1}", index=False, header=False)
            else:
                # Fallback to pure text extraction if no structural grid table found
                df = pd.DataFrame([["No structural tables found in PDF"]])
                df.to_excel(output_path, index=False, header=False)
            return FileResponse(output_path, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", filename=output_path)
            
        elif target_format.lower() == "ppt":
            # Lightweight dynamic compilation structure for presentation layers
            from pptx import Presentation
            from pptx.util import Inches
            import pdfplumber
            output_path = f"{base_name}.pptx"
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]
            
            with pdfplumber.open(input_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    slide = prs.slides.add_slide(blank_slide_layout)
                    if text:
                        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                        tf = txBox.text_frame
                        tf.text = text
            prs.save(output_path)
            return FileResponse(output_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=output_path)
            
        else:
            raise HTTPException(status_code=400, detail="Invalid target office format. Choose 'word', 'excel', or 'ppt'.")
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {str(e)}")
    finally:
        if os.path.exists(input_path): os.remove(input_path)
