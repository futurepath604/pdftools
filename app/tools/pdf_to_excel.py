import io
import os
import zipfile
from flask import Flask, request, render_template_string, send_file, jsonify

import fitz  # PyMuPDF
from pypdf import PdfReader, PdfWriter
from PIL import Image

# --- ADVANCED LAYOUT PARSERS ---
import pdfplumber
from openpyxl import Workbook as ExcelWorkbook
from pptx import Presentation as PptPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

app = Flask(__name__)

# --- PREMIUM RESPONSIVE UI CONTAINER (ALL DEVICE COMPATIBLE) ---
HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>FILEORA - Smart Tools for Every File</title>
<style>
  :root {
    --paper: #F8F9FA;
    --paper-dim: #E9ECEF;
    --ink: #212529;
    --ink-soft: #495057;
    --rule: #CED4DA;
    --accent: #A8431F;
    --accent-deep: #7C2F14;
    --stamp: #2E4D3F;
    --sidebar-w: 280px;
  }
  * { box-sizing: border-box; margin: 0; padding: 0; }
  html, body { height: 100%; font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif; background: var(--paper); color: var(--ink); }
  
  header.masthead {
    display: flex; justify-content: space-between; align-items: center;
    border-bottom: 2px solid var(--ink); padding: 12px 20px; background: #fff;
    position: fixed; top: 0; left: 0; right: 0; height: 65px; z-index: 100;
  }
  header.masthead .brand { font-size: 18px; font-weight: 700; }
  header.masthead .brand small { display: block; font-family: monospace; font-size: 10px; color: var(--accent); letter-spacing: 1px; text-transform: uppercase; }
  
  .menu-toggle {
    display: none; background: var(--accent); color: #fff; border: none;
    padding: 8px 12px; border-radius: 4px; cursor: pointer; font-weight: 600; font-size: 13px;
  }

  .layout-container { display: flex; height: 100%; position: relative; top: 65px; box-sizing: border-box; }
  
  aside.sidebar {
    width: var(--sidebar-w); background: #1A1C1E; border-right: 1px solid var(--rule);
    display: flex; flex-direction: column; overflow-y: auto; height: calc(100vh - 65px);
    transition: transform 0.3s ease; z-index: 90;
  }
  .sidebar-title { color: #8E959E; font-size: 10px; font-weight: 700; text-transform: uppercase; letter-spacing: 1px; padding: 15px 20px 5px; }
  .menu-btn {
    background: none; border: none; color: #DEE2E6; font-size: 13px; text-align: left;
    padding: 12px 20px; cursor: pointer; display: flex; align-items: center; gap: 10px; width: 100%; border-bottom: 1px solid #2D3135;
  }
  .menu-btn:hover, .menu-btn.active { background: #2D3135; color: #fff; font-weight: 600; }
  .icon-span { font-size: 16px; min-width: 20px; text-align: center; }
  
  main.workspace { flex: 1; padding: 30px; overflow-y: auto; height: calc(100vh - 65px); background: var(--paper); }
  .welcome-box { max-width: 900px; margin: 0 auto; text-align: center; padding-top: 10px; }
  
  .grid { display: grid; grid-template-columns: repeat(auto-fill, minmax(240px, 1fr)); gap: 16px; text-align: left; }
  .card { background: #fff; border: 1px solid var(--rule); border-radius: 6px; padding: 16px; cursor: pointer; display: flex; flex-direction: column; gap: 6px; transition: transform 0.2s; }
  .card:hover { border-color: var(--accent); transform: translateY(-2px); box-shadow: 0 6px 12px rgba(0,0,0,0.04); }
  .card-title { font-size: 14px; font-weight: 600; display: flex; align-items: center; gap: 8px; }
  .card p { margin: 0; font-size: 12px; color: var(--ink-soft); line-height: 1.4; }

  .panel { display: none; max-width: 760px; margin: 0 auto; background: #fff; border: 1px solid var(--rule); padding: 25px; border-radius: 6px; box-shadow: 0 4px 12px rgba(0,0,0,0.02); }
  .panel.active { display: block; }
  .panel h2 { font-size: 18px; margin-bottom: 6px; display: flex; align-items: center; gap: 10px; }
  .panel .sub { color: var(--ink-soft); font-size: 13px; margin-bottom: 20px; }
  
  .nav-action-links { display: flex; gap: 15px; margin-bottom: 20px; }
  .action-link { font-size: 12.5px; color: var(--accent); cursor: pointer; text-decoration: underline; font-weight: 600; background: none; border: none; }
  .action-link.home-link { color: var(--ink-soft); }

  .dropzone { border: 2px dashed var(--rule); background: var(--paper); padding: 35px 20px; text-align: center; position: relative; cursor: pointer; margin-bottom: 16px; border-radius: 4px; }
  .dropzone input[type=file] { position: absolute; inset: 0; opacity: 0; cursor: pointer; width: 100%; height: 100%; }
  
  .queue-box { background: #fff; border: 1px solid var(--rule); border-radius: 4px; margin-bottom: 20px; display: none; padding: 14px; }
  .queue-title { font-size: 12px; font-weight: 600; margin-bottom: 10px; color: var(--ink); text-transform: uppercase; }
  .queue-list { display: flex; flex-direction: column; gap: 8px; }
  .queue-item { display: flex; align-items: center; justify-content: space-between; background: var(--paper); padding: 8px 12px; border-radius: 4px; border: 1px solid var(--rule); font-size: 12.5px; }
  .queue-item .remove-node { color: var(--accent); font-weight: bold; cursor: pointer; border: none; background: none; }

  .extra-field { margin-bottom: 16px; }
  .extra-field label { display: block; font-size: 13px; font-weight: 600; margin-bottom: 6px; color: var(--ink-soft); }
  .extra-field input, .extra-field select { width: 100%; padding: 10px; border: 1px solid var(--rule); border-radius: 4px; font-size: 13px; background: #fff; }
  
  .save-info-box { background: #EDF2F7; border: 1px solid #CBD5E0; border-radius: 4px; padding: 12px; margin-bottom: 16px; font-size: 12px; color: #4A5568; line-height: 1.4; display: flex; align-items: center; gap: 8px; }

  button.run { background: var(--accent); color: #fff; border: none; padding: 12px 20px; font-size: 13px; font-weight: 600; border-radius: 4px; cursor: pointer; margin-top: 10px; width: 100%; }
  button.run:hover { background: var(--accent-deep); }
  button.run:disabled { background: var(--rule); cursor: not-allowed; }
  
  .status { margin-top: 16px; font-size: 13px; font-weight: 500; color: var(--stamp); }
  .status.err { color: var(--accent); }

  @media (max-width: 991px) {
    aside.sidebar { position: absolute; left: 0; top: 0; bottom: 0; transform: translateX(-100%); }
    aside.sidebar.open { transform: translateX(0); }
    .menu-toggle { display: block; }
    main.workspace { padding: 20px; }
  }
  @media (max-width: 576px) {
    header.masthead .brand { font-size: 16px; }
    main.workspace { padding: 15px; }
    .panel { padding: 16px; }
    .grid { grid-template-columns: 1fr; }
  }
</style>
</head>
<body>

<header class="masthead">
  <div class="brand">FILEORA<small>Smart Tools for Every File</small></div>
  <button class="menu-toggle" onclick="toggleMobileSidebar()">Tools Menu</button>
</header>

<div class="layout-container">
  <aside class="sidebar" id="app-sidebar">
    <div class="sidebar-title">Main Dashboard</div>
    <button class="menu-btn active" id="btn-menu-home" onclick="resetToHome()">Workspace Arena</button>
    <div class="sidebar-title">Standard PDF & Image Tools</div>
    <div id="sidebar-standard"></div>
  </aside>

  <main class="workspace">
    <div id="home-view" class="welcome-box" style="padding-top: 0; margin-bottom: 20px;">
      <p style="font-size: 14px; color: var(--ink-soft); line-height: 1.5; max-width: 750px; margin: 0 auto 15px auto; font-weight: 500; text-align: center;">
        FILEORA provides smart online tools to convert, compress, resize, and optimize PDF and image files quickly and securely. Simple, fast, and free file solutions for everyone.
      </p>
      <div class="grid" id="grid-all-cards"></div>
    </div>
    <div id="dynamic-panels"></div>
  </main>
</div>

<script>
const TOOLS = [
  {id:"merge_pdf", icon:"📄", name:"Merge PDF", desc:"Combine multiple PDF documents together. Adjust sequence queue arrays dynamically.", action:"/merge-pdf", btnText:"Combine & Download Processed File", accept:"application/pdf", multiple:true},
  {id:"split_pdf_direct", icon:"✂️", name:"Split PDF (Direct)", desc:"Directly split every single page of your source PDF into separate files packaged inside a ZIP container.", action:"/split-pdf-direct", btnText:"Split to Pages & Download ZIP", accept:"application/pdf", multiple:false},
  {id:"split_pdf", icon:"📄", name:"Split PDF (Custom Range)", desc:"Extract specific pages or custom ranges into separate PDF documents.", action:"/split-pdf", btnText:"Split Selected Range & Download", accept:"application/pdf", extra:[{id:"split_range", name:"Page Range to Extract (e.g., 1-3, 5)", type:"text", val:"1-3"}]},
  {id:"rotate_pdf", icon:"📄", name:"Rotate PDF Pages", desc:"Rotate the pages of multiple PDF documents with custom directions.", action:"/rotate-pdf", btnText:"Apply Rotation & Download", accept:"application/pdf", multiple:true, extra:[{id:"rot_deg", name:"Rotation Configuration", type:"select", options:[
    {v:"90",t:"90° Clockwise"},{v:"180",t:"180° Clockwise"},{v:"270",t:"270° Clockwise"},{v:"360",t:"360° Clockwise"},
    {v:"-90",t:"90° Anti-Clockwise"},{v:"-180",t:"180° Anti-Clockwise"},{v:"-270",t:"270° Anti-Clockwise"},{v:"-360",t:"360° Anti-Clockwise"}
  ]}]},
  {id:"organize_pdf", icon:"📄", name:"Organize Pages", desc:"Rearrange layout ordering, drop specified indices, or append clean blank pages anywhere.", action:"/organize-pdf", btnText:"Rebuild Map & Download File", accept:"application/pdf", extra:[{id:"org_map", name:"Target Sequence Blueprint Map (e.g., 1, blank, 2)", type:"text", val:"1, blank, 2"}]},
  {id:"delete_pages", icon:"📄", name:"Delete Pages", desc:"Purge targeted unwanted page streams completely from core file.", action:"/delete-pages", btnText:"Slice & Download Clean Document", accept:"application/pdf", extra:[{id:"del_pages", name:"Pages to drop (e.g., 2, 4-6)", type:"text", val:"2"}]},
  {id:"create_pdf", icon:"📄", name:"Create PDF Engine", desc:"Compile high-fidelity PDF documents directly from images or raw graphic pipelines.", action:"/create-pdf", btnText:"Compile Images & Download PDF", accept:"image/*", multiple:true},
  {id:"images_to_pdf_ext", icon:"📄", name:"Images to PDF (All Formats)", desc:"Convert any asset format bundle seamlessly into a unified master PDF framework.", action:"/images-to-pdf-extended", btnText:"Assemble & Download Framework", accept:"image/*", multiple:true},
  {id:"compress_pdf", icon:"📄", name:"Compress PDF Size", desc:"Shrink asset sizes via smart variable garbage mapping policies.", action:"/compress-pdf", btnText:"Compress & Download Document", accept:"application/pdf", extra:[{id:"comp_level", name:"Compression Strategy", type:"select", options:[{v:"medium",t:"Recommended Compression"},{v:"high",t:"Extreme Compression"},{v:"low",t:"Low Compression"}]}]},
  {id:"compress_image", icon:"📄", name:"Compress Image Size", desc:"Reduce graphic image payload with custom adjusted dimensions and target resolution rules.", action:"/compress-image", btnText:"Optimize & Download Images", accept:"image/*", multiple:true, extra:[{id:"img_quality", name:"Compression Quality (1-100)", type:"text", val:"75"},{id:"img_scale", name:"Scale Resolution Width Factor (e.g., 100% means unchanged)", type:"select", options:[{v:"100",t:"100% Original Resolution"},{v:"75",t:"75% Scale"},{v:"50",t:"50% Compact Resolution"},{v:"25",t:"25% Thumbnail Size"}]}]},
  {id:"rotate_image", icon:"📄", name:"Rotate Image", desc:"Rotate multiple images concurrently with flexible choice of degree, direction, and save file format.", action:"/rotate-image", btnText:"Execute Rotation & Download Image", accept:"image/*", multiple:true, extra:[
    {id:"img_rot_deg", name:"Rotation Angle & Direction", type:"select", options:[
      {v:"90",t:"90° Clockwise"},{v:"180",t:"180° Clockwise"},{v:"270",t:"270° Clockwise"},{v:"360",t:"360° Clockwise"},
      {v:"-90",t:"90° Anti-Clockwise"},{v:"-180",t:"180° Anti-Clockwise"},{v:"-270",t:"270° Anti-Clockwise"},{v:"-360",t:"360° Anti-Clockwise"}
    ]},
    {id:"img_out_fmt", name:"Save File Option Target Format", type:"select", options:[
      {v:"original",t:"Keep Original Format"},{v:"jpeg",t:"JPEG (.jpg)"},{v:"png",t:"PNG (.png)"},{v:"webp",t:"WEBP (.webp)"},{v:"bmp",t:"BMP (.bmp)"},{v:"tiff",t:"TIFF (.tiff)"}
    ]}
  ]},
  {id:"lock_pdf_secure", icon:"📄", name:"Lock PDF (Password)", desc:"Encrypt file stream containers with master standard key sets to prevent illegal visibility.", action:"/lock-pdf", btnText:"Deploy Keys & Download Locked PDF", accept:"application/pdf", extra:[{id:"lock_pass", name:"Configure Target Key Password", type:"text", val:"SecurePass123"}]},
  {id:"unlock_pdf", icon:"📄", name:"Unlock PDF Password", desc:"Strip restrictions from encrypted layout containers natively.", action:"/unlock-pdf", btnText:"Strip Security Lock & Download", accept:"application/pdf", extra:[{id:"upass", name:"Current Password", type:"text", val:""}]},
  {id:"pdf2img", icon:"📄", name:"PDF to Images", desc:"Extract document elements into system supported extension frameworks.", action:"/pdf-to-jpg", btnText:"Render Pages & Download ZIP", accept:"application/pdf", extra:[{id:"img_fmt", name:"Target Format", type:"select", options:[{v:"jpg",t:"JPEG (.jpg)"},{v:"png",t:"PNG (.png)"},{v:"webp",t:"WEBP (.webp)"},{v:"tiff",t:"TIFF (.tiff)"}]}]},
  {id:"pdf_to_word_ext", icon:"📄", name:"PDF to Word Converter", desc:"Parse architectural structures into Microsoft Word formats natively.", action:"/pdf-to-word-extended", btnText:"Extract & Download Word Document", accept:"application/pdf", extra:[{id:"word_fmt", name:"Output File Format Extension", type:"select", options:[{v:"docx",t:"Word Document (.docx)"},{v:"doc",t:"Legacy Word Document (.doc)"}]}]},
  {id:"pdf_to_excel_ext", icon:"📊", name:"PDF to Excel Converter", desc:"Extract structural cell schemas and data grids precisely into native spreadsheet format tables.", action:"/pdf-to-excel-extended", btnText:"Compile Spreadsheet & Download", accept:"application/pdf", extra:[{id:"excel_fmt", name:"Output Spreadsheet Extension", type:"select", options:[{v:"xlsx",t:"Excel Workbook (.xlsx)"}]}]},
  {id:"pdf_to_ppt_ext", icon:"📄", name:"PDF to PowerPoint Converter", desc:"Transform static presentation visual sheets into slide format patterns layout.", action:"/pdf-to-ppt-extended", btnText:"Build Slide Compilation & Download", accept:"application/pdf", extra:[{id:"ppt_fmt", name:"Output Presentation Extension", type:"select", options:[{v:"pptx",t:"PowerPoint Presentation (.pptx)"},{v:"ppt",t:"Legacy PowerPoint Layout (.ppt)"}]}
]}
];

let fileRegistry = {};

TOOLS.forEach(t => {
  fileRegistry[t.id] = [];
  const navBtn = document.createElement('button');
  navBtn.className = 'menu-btn'; navBtn.id = `nav-btn-${t.id}`;
  navBtn.innerHTML = `<span class="icon-span">${t.icon}</span> ${t.name}`;
  navBtn.onclick = () => activateToolPanel(t.id);
  document.getElementById('sidebar-standard').appendChild(navBtn);

  const card = document.createElement('div');
  card.className = 'card'; card.innerHTML = `<div class="card-title"><span>${t.icon}</span>${t.name}</div><p>${t.desc}</p>`;
  card.onclick = () => activateToolPanel(t.id);
  document.getElementById('grid-all-cards').appendChild(card);

  const panel = document.createElement('section');
  panel.className = 'panel'; panel.id = 'panel-' + t.id;
  
  let extrasHtml = '';
  if (t.extra) {
    t.extra.forEach(ex => {
      extrasHtml += `<div class="extra-field"><label>${ex.name}</label>`;
      if (ex.type === 'select') {
        extrasHtml += `<select id="ex-${t.id}-${ex.id}">`;
        ex.options.forEach(o => extrasHtml += `<option value="${o.v}">${o.t}</option>`);
        extrasHtml += `</select>`;
      } else { extrasHtml += `<input type="text" id="ex-${t.id}-${ex.id}" value="${ex.val}">`; }
      extrasHtml += `</div>`;
    });
  }

  panel.innerHTML = `
    <div class="nav-action-links">
      <button class="action-link home-link" onclick="resetToHome()">Return to Home Dashboard</button>
      <button class="action-link" onclick="resetToolMatrix('${t.id}')">Reset Tool Matrix</button>
    </div>
    <h2><span>${t.icon}</span> ${t.name}</h2>
    <div class="sub">${t.desc}</div>
    <div class="dropzone">
      <p id="label-${t.id}">Drop files here, or <span style="color:var(--accent); font-weight:600;">browse storage system</span></p>
      <input type="file" id="input-${t.id}" accept="${t.accept}" ${t.multiple ? 'multiple' : ''}>
    </div>
    <div class="queue-box" id="queue-box-${t.id}">
      <div class="queue-title">Active Target Queue</div>
      <div class="queue-list" id="queue-list-${t.id}"></div>
    </div>
    ${extrasHtml}
    <div class="save-info-box">
      <span>💾 <strong>Save & Download Location Rule:</strong> Files are securely rendered on the fly and immediately pushed as a direct download. Your browser will prompt for directory location choice or save it into your default system storage folder.</span>
    </div>
    <button class="run" id="btn-${t.id}" disabled>${t.btnText}</button>
    <div class="status" id="status-${t.id}"></div>
  `;
  document.getElementById('dynamic-panels').appendChild(panel);

  const input = document.getElementById(`input-${t.id}`);
  input.addEventListener('change', (e) => {
    const incoming = Array.from(e.target.files);
    fileRegistry[t.id] = t.multiple ? fileRegistry[t.id].concat(incoming) : incoming;
    renderQueueList(t.id);
  });

  const runBtn = document.getElementById(`btn-${t.id}`);
  runBtn.onclick = async () => {
    const formData = new FormData();
    const activeFiles = fileRegistry[t.id];
    if(activeFiles.length === 0) return;
    activeFiles.forEach(f => formData.append('files', f));

    let selectedExtension = "pdf";
    if (t.extra) {
      t.extra.forEach(ex => {
        const val = document.getElementById(`ex-${t.id}-${ex.id}`).value;
        formData.append(ex.id, val);
        if (ex.id.endsWith('_fmt')) { selectedExtension = val; }
      });
    }

    const statusObj = document.getElementById(`status-${t.id}`);
    statusObj.textContent = 'Processing high-fidelity matrix render pipelines & downloading...';
    statusObj.classList.remove('err');
    runBtn.disabled = true;

    try {
      const response = await fetch(t.action, { method: 'POST', body: formData });
      
      if (!response.ok) {
        let errMsg = "Execution failed inside core process.";
        try {
          const errData = await response.json();
          errMsg = errData.error || errMsg;
        } catch(pErr) {
          errMsg = `Server Error HTTP: ${response.status}`;
        }
        throw new Error(errMsg);
      }

      const blob = await response.blob();
      const url = window.URL.createObjectURL(blob);
      const a = document.createElement('a'); a.href = url;
      
      let ext = "_" + selectedExtension;
      if (t.id === "lock_pdf_secure") ext = "_locked.pdf";
      if (t.id === "split_pdf_direct") ext = "_pages_package.zip";
      if (t.id === "compress_image" && activeFiles.length === 1) ext = "_optimized." + activeFiles[0].name.split('.').pop();
      if (t.id === "rotate_image" && activeFiles.length === 1) {
        const outFmtSelection = document.getElementById(`ex-rotate_image-img_out_fmt`).value;
        const fallbackExt = activeFiles[0].name.split('.').pop();
        ext = "_rotated." + (outFmtSelection === "original" ? fallbackExt : outFmtSelection);
      }
      if (response.headers.get('Content-Type') === 'application/zip') ext = "_package.zip";

      a.download = activeFiles[0].name.split('.')[0] + ext;
      document.body.appendChild(a); a.click(); window.URL.revokeObjectURL(url);
      statusObj.textContent = 'Operations completed successfully. File downloaded.';
    } catch (err) {
      statusObj.textContent = err.message; statusObj.classList.add('err');
    } finally { runBtn.disabled = false; }
  };
});

function renderQueueList(toolId) {
  const listObj = document.getElementById(`queue-list-${toolId}`);
  const containerBox = document.getElementById(`queue-box-${toolId}`);
  const labelObj = document.getElementById(`label-${toolId}`);
  const executeBtn = document.getElementById(`btn-${toolId}`);
  const arr = fileRegistry[toolId];

  listObj.innerHTML = '';
  if(arr.length === 0) {
    if(containerBox) containerBox.style.display = 'none';
    labelObj.textContent = "Drop files here, or browse storage system";
    if(executeBtn) executeBtn.disabled = true;
    return;
  }
  if(containerBox) containerBox.style.display = 'block';
  labelObj.textContent = `${arr.length} item streams loaded successfully.`;
  if(executeBtn) executeBtn.disabled = false;

  arr.forEach((file, index) => {
    const row = document.createElement('div'); row.className = 'queue-item';
    row.innerHTML = `<span><strong>[${index + 1}]</strong> ${file.name}</span><button class="remove-node" onclick="removeQueueNode('${toolId}', ${index})">Remove</button>`;
    listObj.appendChild(row);
  });
}

function removeQueueNode(toolId, index) { fileRegistry[toolId].splice(index, 1); renderQueueList(toolId); }
function resetToolMatrix(toolId) {
  fileRegistry[toolId] = [];
  const inputElement = document.getElementById(`input-${toolId}`); if(inputElement) inputElement.value = '';
  const statusObj = document.getElementById(`status-${toolId}`); if(statusObj) { statusObj.textContent = ''; statusObj.classList.remove('err'); }
  renderQueueList(toolId);
}
function toggleMobileSidebar() { document.getElementById('app-sidebar').classList.toggle('open'); }
function resetToHome() {
  document.getElementById('home-view').style.display = 'block';
  document.querySelectorAll('.panel').forEach(p => p.classList.remove('active'));
  document.querySelectorAll('.menu-btn').forEach(b => b.classList.remove('active'));
  document.getElementById('btn-menu-home').classList.add('active');
  document.getElementById('app-sidebar').classList.remove('open');
}
function activateToolPanel(id) {
  document.getElementById('home-view').style.display = 'none';
  document.querySelectorAll('.panel').forEach(p => p.classList.remove('active'));
  document.querySelectorAll('.menu-btn').forEach(b => b.classList.remove('active'));
  document.getElementById('panel-' + id).classList.add('active');
  document.getElementById('nav-btn-' + id).classList.add('active');
  document.getElementById('app-sidebar').classList.remove('open');
}
</script>
</body>
</html>
"""

@app.route('/')
def home(): return render_template_string(HTML_TEMPLATE)

@app.route('/merge-pdf', methods=['POST'])
def merge_pdf():
    files = request.files.getlist('files')
    try:
        writer = PdfWriter()
        for f in files:
            reader = PdfReader(io.BytesIO(f.read()))
            for page in reader.pages: writer.add_page(page)
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype="application/pdf", as_attachment=True, download_name="combined.pdf")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/split-pdf-direct', methods=['POST'])
def split_pdf_direct():
    files = request.files.getlist('files')
    if not files: return jsonify({"error": "No file uploaded"}), 400
    try:
        file = files[0]
        reader = PdfReader(io.BytesIO(file.read()))
        total_pages = len(reader.pages)
        
        zip_stream = io.BytesIO()
        with zipfile.ZipFile(zip_stream, 'w') as zip_file:
            for idx in range(total_pages):
                writer = PdfWriter()
                writer.add_page(reader.pages[idx])
                page_out = io.BytesIO()
                writer.write(page_out)
                zip_file.writestr(f"page_{idx + 1}.pdf", page_out.getvalue())
                
        zip_stream.seek(0)
        return send_file(zip_stream, mimetype="application/zip", as_attachment=True, download_name="split_pages.zip")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/split-pdf', methods=['POST'])
def split_pdf():
    file = request.files.getlist('files')[0]
    ranges = request.form.get('split_range', '1').strip()
    try:
        reader = PdfReader(io.BytesIO(file.read()))
        writer = PdfWriter()
        total_p = len(reader.pages)
        
        for part in ranges.split(','):
            part = part.strip()
            if '-' in part:
                s, e = map(int, part.split('-'))
                for p_idx in range(max(1, s), min(e, total_p) + 1):
                    writer.add_page(reader.pages[p_idx - 1])
            elif part.isdigit():
                idx = int(part)
                if 1 <= idx <= total_p:
                    writer.add_page(reader.pages[idx - 1])
                    
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype="application/pdf", as_attachment=True, download_name="split_output.pdf")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/rotate-pdf', methods=['POST'])
def rotate_pdf():
    files = request.files.getlist('files')
    raw_deg = int(request.form.get('rot_deg', '90'))
    deg = raw_deg if raw_deg >= 0 else (360 + raw_deg)
    
    try:
        if len(files) == 1:
            reader = PdfReader(io.BytesIO(files[0].read()))
            writer = PdfWriter()
            for page in reader.pages:
                page.rotate(deg)
                writer.add_page(page)
            out = io.BytesIO(); writer.write(out); out.seek(0)
            return send_file(out, mimetype="application/pdf", as_attachment=True, download_name="rotated_document.pdf")
        else:
            zip_stream = io.BytesIO()
            with zipfile.ZipFile(zip_stream, 'w') as zip_file:
                for idx, f in enumerate(files):
                    reader = PdfReader(io.BytesIO(f.read()))
                    writer = PdfWriter()
                    for page in reader.pages:
                        page.rotate(deg)
                        writer.add_page(page)
                    out_f = io.BytesIO(); writer.write(out_f)
                    zip_file.writestr(f"rotated_{idx+1}_{f.name}", out_f.getvalue())
            zip_stream.seek(0)
            return send_file(zip_stream, mimetype="application/zip", as_attachment=True, download_name="rotated_bundle.zip")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/organize-pdf', methods=['POST'])
def organize_pdf():
    files = request.files.getlist('files')
    map_str = request.form.get('org_map', '1, blank, 2')
    try:
        writer = PdfWriter(); all_pages = []
        for f in files: all_pages.extend(PdfReader(io.BytesIO(f.read())).pages)
        for inst in [x.strip().lower() for x in map_str.split(',')]:
            if inst == 'blank': writer.add_blank_page(width=612, height=792)
            else:
                try:
                    idx = int(inst) - 1
                    if 0 <= idx < len(all_pages): writer.add_page(all_pages[idx])
                except: continue
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype="application/pdf", as_attachment=True, download_name="organized.pdf")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/delete-pages', methods=['POST'])
def delete_pages():
    file = request.files.getlist('files')[0]
    del_str = request.form.get('del_pages', '')
    try:
        reader = PdfReader(io.BytesIO(file.read())); writer = PdfWriter(); drop = set()
        for part in del_str.split(','):
            part = part.strip()
            if '-' in part:
                s, e = map(int, part.split('-'))
                drop.update(range(s-1, e))
            elif part.isdigit(): drop.add(int(part)-1)
        for i in range(len(reader.pages)):
            if i not in drop: writer.add_page(reader.pages[i])
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype="application/pdf", as_attachment=True, download_name="purged.pdf")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/create-pdf', methods=['POST'])
def create_pdf():
    files = request.files.getlist('files')
    try:
        doc = fitz.open()
        for f in files:
            img_doc = fitz.open(stream=f.read(), filetype=f.name.split('.')[-1])
            pdf_bytes = img_doc.convert_to_pdf(); img_doc.close()
            p_stream = fitz.open("pdf", pdf_bytes); doc.insert_pdf(p_stream); p_stream.close()
        out = io.BytesIO(doc.write()); doc.close(); out.seek(0)
        return send_file(out, mimetype="application/pdf", as_attachment=True, download_name="compiled.pdf")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/images-to-pdf-extended', methods=['POST'])
def images_to_pdf_extended(): return create_pdf()

@app.route('/compress-pdf', methods=['POST'])
def compress_pdf():
    file = request.files.getlist('files')[0]
    level = request.form.get('comp_level', 'medium')
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        out_bytes = doc.write(garbage=4, deflate=True, clean=(1 if level in ['high','medium'] else 0))
        doc.close()
        return send_file(io.BytesIO(out_bytes), mimetype="application/pdf", as_attachment=True, download_name="compressed.pdf")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/compress-image', methods=['POST'])
def compress_image():
    files = request.files.getlist('files')
    quality = int(request.form.get('img_quality', '75'))
    scale = int(request.form.get('img_scale', '100'))
    try:
        if len(files) == 1:
            img = Image.open(io.BytesIO(files[0].read()))
            orig_format = img.format or "JPEG"
            
            if scale != 100:
                nw = int(img.width * (scale / 100.0))
                nh = int(img.height * (scale / 100.0))
                img = img.resize((nw, nh), Image.Resampling.LANCZOS)
                
            if img.mode in ('RGBA', 'P') and orig_format in ('JPEG', 'JPG'):
                img = img.convert('RGB')
                
            out = io.BytesIO()
            img.save(out, format=orig_format, quality=quality, optimize=True)
            out.seek(0)
            
            mimetype = f"image/{orig_format.lower()}"
            return send_file(out, mimetype=mimetype, as_attachment=True, download_name=f"optimized.{orig_format.lower()}")
        else:
            zip_stream = io.BytesIO()
            with zipfile.ZipFile(zip_stream, 'w') as zip_file:
                for idx, f in enumerate(files):
                    img = Image.open(io.BytesIO(f.read()))
                    orig_format = img.format or "JPEG"
                    if scale != 100:
                        img = img.resize((int(img.width * (scale/100.0)), int(img.height * (scale/100.0))), Image.Resampling.LANCZOS)
                    if img.mode in ('RGBA', 'P') and orig_format in ('JPEG', 'JPG'):
                        img = img.convert('RGB')
                    out_f = io.BytesIO()
                    img.save(out_f, format=orig_format, quality=quality, optimize=True)
                    zip_file.writestr(f"optimized_{idx+1}.{orig_format.lower()}", out_f.getvalue())
            zip_stream.seek(0)
            return send_file(zip_stream, mimetype="application/zip", as_attachment=True, download_name="optimized_images.zip")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/rotate-image', methods=['POST'])
def rotate_image():
    files = request.files.getlist('files')
    raw_deg = int(request.form.get('img_rot_deg', '90'))
    out_fmt = request.form.get('img_out_fmt', 'original')
    deg = raw_deg if raw_deg >= 0 else (360 + raw_deg)
    
    try:
        if len(files) == 1:
            f = files[0]
            img = Image.open(io.BytesIO(f.read()))
            orig_ext = f.name.split('.')[-1].upper()
            orig_format = img.format or (orig_ext if orig_ext in ['JPEG', 'PNG', 'WEBP', 'BMP', 'TIFF'] else 'JPEG')
            
            # PIL rotate is counter-clockwise for positive values, map correctly
            pil_deg = (360 - deg) % 360
            img = img.rotate(pil_deg, expand=True)
            
            target_format = orig_format if out_fmt == 'original' else out_fmt.upper()
            if target_format == 'JPG': target_format = 'JPEG'
            
            if img.mode in ('RGBA', 'P') and target_format == 'JPEG':
                img = img.convert('RGB')
                
            out = io.BytesIO()
            img.save(out, format=target_format)
            out.seek(0)
            
            ext_map = {'JPEG': 'jpg', 'PNG': 'png', 'WEBP': 'webp', 'BMP': 'bmp', 'TIFF': 'tiff'}
            final_ext = ext_map.get(target_format, 'jpg')
            return send_file(out, mimetype=f"image/{final_ext}", as_attachment=True, download_name=f"rotated.{final_ext}")
        else:
            zip_stream = io.BytesIO()
            with zipfile.ZipFile(zip_stream, 'w') as zip_file:
                for idx, f in enumerate(files):
                    img = Image.open(io.BytesIO(f.read()))
                    orig_ext = f.name.split('.')[-1].upper()
                    orig_format = img.format or (orig_ext if orig_ext in ['JPEG', 'PNG', 'WEBP', 'BMP', 'TIFF'] else 'JPEG')
                    
                    pil_deg = (360 - deg) % 360
                    img = img.rotate(pil_deg, expand=True)
                    
                    target_format = orig_format if out_fmt == 'original' else out_fmt.upper()
                    if target_format == 'JPG': target_format = 'JPEG'
                    
                    if img.mode in ('RGBA', 'P') and target_format == 'JPEG':
                        img = img.convert('RGB')
                        
                    out_f = io.BytesIO()
                    img.save(out_f, format=target_format)
                    
                    ext_map = {'JPEG': 'jpg', 'PNG': 'png', 'WEBP': 'webp', 'BMP': 'bmp', 'TIFF': 'tiff'}
                    final_ext = ext_map.get(target_format, 'jpg')
                    zip_file.writestr(f"rotated_{idx+1}.{final_ext}", out_f.getvalue())
            zip_stream.seek(0)
            return send_file(zip_stream, mimetype="application/zip", as_attachment=True, download_name="rotated_images.zip")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/lock-pdf', methods=['POST'])
def lock_pdf():
    file = request.files.getlist('files')[0]
    password = request.form.get('lock_pass', 'SecurePass123')
    try:
        reader = PdfReader(io.BytesIO(file.read()))
        writer = PdfWriter()
        for page in reader.pages: writer.add_page(page)
        writer.encrypt(password)
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype="application/pdf", as_attachment=True, download_name="locked.pdf")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/unlock-pdf', methods=['POST'])
def unlock_pdf():
    file = request.files.getlist('files')[0]
    password = request.form.get('upass', '')
    try:
        reader = PdfReader(io.BytesIO(file.read()))
        if reader.is_encrypted:
            reader.decrypt(password)
        writer = PdfWriter()
        for page in reader.pages: writer.add_page(page)
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype="application/pdf", as_attachment=True, download_name="unlocked.pdf")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/pdf-to-jpg', methods=['POST'])
def pdf_to_jpg():
    file = request.files.getlist('files')[0]
    img_fmt = request.form.get('img_fmt', 'jpg').lower()
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        zip_stream = io.BytesIO()
        with zipfile.ZipFile(zip_stream, 'w') as zip_file:
            for idx, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes(img_fmt if img_fmt != 'jpg' else 'jpeg')
                ext = 'jpg' if img_fmt == 'jpg' else img_fmt
                zip_file.writestr(f"page_{idx + 1}.{ext}", img_bytes)
        doc.close()
        zip_stream.seek(0)
        return send_file(zip_stream, mimetype="application/zip", as_attachment=True, download_name="pdf_images.zip")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/pdf-to-word-extended', methods=['POST'])
def pdf_to_word_extended():
    file = request.files.getlist('files')[0]
    try:
        # Simple text extraction fallback to docx if pdf2docx has complex environment dependencies
        from docx import Document
        doc = Document()
        reader = PdfReader(io.BytesIO(file.read()))
        for idx, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                for line in text.split('\n'):
                    doc.add_paragraph(line)
            if idx < len(reader.pages) - 1:
                doc.add_page_break()
        out = io.BytesIO()
        doc.save(out)
        out.seek(0)
        return send_file(out, mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document", as_attachment=True, download_name="converted_document.docx")
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/pdf-to-excel-extended', methods=['POST'])
def pdf_to_excel_extended():
    files = request.files.getlist('files')
    if not files:
        return jsonify({"error": "No file uploaded"}), 400
    try:
        file = files[0]
        pdf_bytes = file.read()
        
        wb = ExcelWorkbook()
        default_sheet = wb.active
        wb.remove(default_sheet)
        
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for page_idx, page in enumerate(pdf.pages):
                ws = wb.create_sheet(title=f"Page_{page_idx + 1}")
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        for row in table:
                            cleaned_row = [cell if cell is not None else "" for cell in row]
                            ws.append(cleaned_row)
                        ws.append([])
                else:
                    text = page.extract_text()
                    if text:
                        for line in text.split('\n'):
                            ws.append([line])
                            
        out = io.BytesIO()
        wb.save(out)
        out.seek(0)
        
        return send_file(
            out, 
            mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", 
            as_attachment=True, 
            download_name="converted_document.xlsx"
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/pdf-to-ppt-extended', methods=['POST'])
def pdf_to_ppt_extended():
    file = request.files.getlist('files')[0]
    try:
        prs = PptPresentation()
        doc = fitz.open(stream=file.read(), filetype="pdf")
        blank_slide_layout = prs.slide_layouts[6]
        for page in doc:
            slide = prs.slides.add_slide(blank_slide_layout)
            text = page.get_text()
            if text:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(6.0))
                tf = txBox.text_frame
                tf.word_wrap = True
                tf.text = text
        doc.close()
        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        return send_file(out, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation", as_attachment=True, download_name="converted_presentation.pptx")
    except Exception as e: return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
