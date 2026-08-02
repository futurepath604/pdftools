import os
import fitz  # PyMuPDF
from fastapi import APIRouter, File, UploadFile, HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Inches

router = APIRouter(prefix="/api/tools", tags=["PDF to PPT"])
TEMP_DIR = "/tmp"
os.makedirs(TEMP_DIR, exist_ok=True)

@router.post("/pdf-to-ppt")
async def convert_pdf_to_ppt(file: UploadFile = File(...)):
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    input_path = os.path.join(TEMP_DIR, file.filename)
    output_filename = f"{os.path.splitext(file.filename)[0]}.pptx"
    output_path = os.path.join(TEMP_DIR, output_filename)
    
    try:
        with open(input_path, "wb") as buffer:
            buffer.write(await file.read())
            
        prs = Presentation()
        doc = fitz.open(input_path)
        
        blank_slide_layout = prs.slide_layouts[6]
        
        for page in doc:
            pix = page.get_pixmap(dpi=150)
            img_path = os.path.join(TEMP_DIR, f"slide_{page.number}.png")
            pix.save(img_path)
            
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            
        doc.close()
        prs.save(output_path)
        
        return FileResponse(output_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=output_filename)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF to PPT conversion failed: {str(e)}")
